from flask import Flask, render_template, request, jsonify, url_for, session, redirect, flash, send_file
from pdf_processor import PDFProcessor
from rag_engine import RAGEngine
from dotenv import load_dotenv
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import re
from pypdf import PdfReader
import base64
import io
import fitz  
import requests
from io import BytesIO
from PIL import Image
import random
from werkzeug.middleware.dispatcher import DispatcherMiddleware
from functools import wraps
import time

load_dotenv()

APPLICATION_ROOT = os.getenv('APPLICATION_ROOT', '')

app = Flask(__name__, static_url_path=APPLICATION_ROOT + '/static')
app.secret_key = os.getenv('SECRET_KEY', 'markaidefaultsecretkey') # Add secret key for sessions

if APPLICATION_ROOT:
    app.config['APPLICATION_ROOT'] = '/marketingai'
    app.wsgi_app = DispatcherMiddleware(Flask('dummy_app'), {
        APPLICATION_ROOT: app.wsgi_app
    })


def get_redirect_url(endpoint):
    if APPLICATION_ROOT:
        return f"{APPLICATION_ROOT}{url_for(endpoint)}"
    else:
        return url_for(endpoint)

pdf_processor = PDFProcessor()

UPLOAD_FOLDER = 'static/uploads'
VECTOR_STORE_PATH = 'static/vector_store'
PROCESSED_FILES_LOG = 'static/processed_files.json'

for path in [UPLOAD_FOLDER, VECTOR_STORE_PATH]:
    if not os.path.exists(path):
        os.makedirs(path)

rag_engine = RAGEngine(vector_store_path=VECTOR_STORE_PATH, processed_files_log=PROCESSED_FILES_LOG)

class PPTGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.title_slide_layout = self.prs.slide_layouts[0]
        self.section_header_layout = self.prs.slide_layouts[2]
        self.title_content_layout = self.prs.slide_layouts[1]
        self.two_content_layout = self.prs.slide_layouts[3]
        self.comparison_layout = self.prs.slide_layouts[4]
        self.blank_layout = self.prs.slide_layouts[6]
        
        self.colors = {
            'primary': RGBColor(0x00, 0x45, 0x89),
            'secondary': RGBColor(0x60, 0x7D, 0x8B),
            'accent1': RGBColor(0x00, 0x96, 0xC7),
            'accent2': RGBColor(0xF2, 0x7D, 0x52),
            'accent3': RGBColor(0xFF, 0xB7, 0x4D),
            'dark': RGBColor(0x33, 0x33, 0x33),
            'text': RGBColor(0x24, 0x29, 0x39),
            'lightText': RGBColor(0x6C, 0x75, 0x7D)
        }
        
        self.shapes = [
            MSO_SHAPE.RECTANGLE,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            MSO_SHAPE.OVAL
        ]
        
        self.unsplash_api_key = os.getenv('UNSPLASH_ACCESS_KEY')
        self.unsplash_url = "https://api.unsplash.com/search/photos"

    def get_image_from_unsplash(self, query):
        try:
            if not self.unsplash_api_key:
                print("Unsplash API key not found")
                return None
                
            search_terms = query.replace("slide", "").replace("presentation", "")
            search_terms = re.sub(r'[^\w\s]', ' ', search_terms)
            search_terms = " ".join(search_terms.split()[:5])
            
            response = requests.get(
                self.unsplash_url,
                params={
                    "query": search_terms,
                    "per_page": 3,
                    "orientation": "landscape",
                    "client_id": self.unsplash_api_key
                }
            )
            
            if response.status_code == 200:
                data = response.json()
                if data["results"]:
                    image_url = random.choice(data["results"])["urls"]["regular"]
                    image_response = requests.get(image_url)
                    if image_response.status_code == 200:
                        return BytesIO(image_response.content)
            
            return None
        except Exception as e:
            print(f"Error fetching image: {str(e)}")
            return None
    
    def apply_text_formatting(self, shape, size=None, color=None, bold=False, alignment=PP_ALIGN.LEFT):
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = alignment
            if size:
                paragraph.font.size = Pt(size)
            if color:
                paragraph.font.color.rgb = color
            paragraph.font.bold = bold
            paragraph.line_spacing = 1.2
        
        return text_frame.paragraphs[0]
    
    def add_background_style(self, slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
        
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.25), self.prs.slide_height
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.colors['accent1']
        accent_bar.line.fill.background()

    def create_slide_title(self, slide, title_text, size=28, color=None, alignment=PP_ALIGN.LEFT):
        title = slide.shapes.title
        title.text = title_text
        if not color:
            color = self.colors['primary']
        self.apply_text_formatting(title, size=size, color=color, bold=True, alignment=alignment)
        
        title_shape = slide.shapes.title
        underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            title_shape.left, title_shape.top + title_shape.height + Pt(2),
            Inches(2), Pt(3)
        )
        underline.fill.solid()
        underline.fill.fore_color.rgb = self.colors['accent1']
        underline.line.fill.background()
    
    def add_image_to_slide(self, slide, image_stream, left=None, top=None, width=None, height=None):
        if not image_stream:
            return None
            
        if left is None:
            if slide.slide_layout == self.title_slide_layout:
                left = Inches(2.5)
            elif slide.slide_layout == self.two_content_layout:
                left = Inches(5)
            else:
                left = Inches(1)
        
        if top is None:
            top = Inches(2.5)
            
        if width is None:
            if slide.slide_layout == self.title_slide_layout:
                width = Inches(5)
            elif slide.slide_layout == self.two_content_layout:
                width = Inches(4)
            else:
                width = Inches(3)
        
        if height is None:
            height = Inches(3)
            
        try:
            image = Image.open(image_stream)
            img_ratio = image.width / image.height
            
            calculated_height = width / img_ratio
            if calculated_height > height:
                width = height * img_ratio
            else:
                height = calculated_height
                
            image_stream.seek(0)
            
            picture = slide.shapes.add_picture(
                image_stream, left, top, width, height
            )
            
            return picture
        except Exception as e:
            print(f"Error adding image: {str(e)}")
            return None
    
    def add_footer(self, slide, text="Generated with AI Assistant"):
        footer = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.9), Inches(9), Inches(0.3)
        )
        footer_p = footer.text_frame.add_paragraph()
        footer_p.text = text
        footer_p.alignment = PP_ALIGN.RIGHT
        footer_p.font.size = Pt(8)
        footer_p.font.color.rgb = self.colors['lightText']
    
    def add_bullet_points(self, shape, points, level_styles=None):
        level_styles = {
            0: {"size": 18, "color": self.colors['text'], "bold": False},
            1: {"size": 16, "color": self.colors['text'], "bold": False},
            2: {"size": 14, "color": self.colors['lightText'], "bold": False}
        }
        
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.auto_size = True
        
        for point in points:
            point = point.strip()
            if not point:
                continue
                
            level = 0
            if point.startswith('  ') or point.startswith('\t'):
                level = 1
                point = point.lstrip()
            if point.startswith('   ') or point.startswith('\t\t'):
                level = 2
                point = point.lstrip()
                
            point = re.sub(r'^[-•*]\s+', '', point)
                
            p = text_frame.add_paragraph()
            p.text = point
            p.level = level
            
            style = level_styles.get(level, level_styles[0])
            p.font.size = Pt(style["size"])
            p.font.color.rgb = style["color"]
            p.font.bold = style["bold"]
            
            p.space_before = Pt(4)
            p.space_after = Pt(4)

    def parse_content(self, content):
        content = re.sub(r'<br\s*/?>', '\n', content)
        content = re.sub(r'<[^>]+>', '', content)
        
        slides = []
        current_slide = {"title": "", "content": [], "type": "content"}
        
        lines = content.split('\n')
        slide_markers = ['slide', 'page', 'section']
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            is_slide_title = False
            line_lower = line.lower()
            
            if any(marker in line_lower for marker in slide_markers) and ':' in line:
                if current_slide["title"]:
                    slides.append(current_slide)
                title_text = line.split(':', 1)[1].strip()
                current_slide = {"title": title_text, "content": [], "type": "content"}
                is_slide_title = True
            
            elif line.isupper() or (re.match(r'^[A-Z][A-Z\s]+$', line)) or line.endswith(':'):
                if current_slide["title"]:
                    slides.append(current_slide)
                title_text = line.rstrip(':')
                current_slide = {"title": title_text, "content": [], "type": "content"}
                is_slide_title = True
                
            if not is_slide_title:
                current_slide["content"].append(line)
                
        if current_slide["title"]:
            slides.append(current_slide)
            
        if slides:
            slides[0]["type"] = "title"
            
            for i, slide in enumerate(slides[1:], 1):
                title_lower = slide["title"].lower()
                if (len(slide["content"]) <= 2 and 
                    any(word in title_lower for word in 
                        ["section", "part", "agenda", "overview", "introduction", "conclusion"])):
                    slide["type"] = "section"
                    
                elif any(word in title_lower for word in 
                         ["vs", "versus", "comparison", "compare", "differences", "similarities"]):
                    slide["type"] = "comparison"
        
        return slides

    def create_presentation(self, content):
        slides_data = self.parse_content(content)
        
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
        
        for i, slide_data in enumerate(slides_data):
            slide_title = slide_data["title"]
            slide_content = slide_data["content"]
            slide_type = slide_data["type"]
            
            if not slide_title and not slide_content:
                continue
            
            if slide_type == "title":
                slide = self.prs.slides.add_slide(self.title_slide_layout)
                self.add_background_style(slide)
                
                title = slide.shapes.title
                title.text = slide_title
                self.apply_text_formatting(title, size=36, color=self.colors['primary'], bold=True, alignment=PP_ALIGN.CENTER)
                
                subtitle = slide.placeholders[1]
                subtitle_text = "Created with AI Marketing Assistant"
                if slide_content:
                    subtitle_text = slide_content[0]
                subtitle.text = subtitle_text
                self.apply_text_formatting(subtitle, size=20, color=self.colors['secondary'], alignment=PP_ALIGN.CENTER)
                
                image_stream = self.get_image_from_unsplash(slide_title)
                if image_stream:
                    picture = self.add_image_to_slide(
                        slide, image_stream, 
                        left=Inches(2.5), top=Inches(1.8), 
                        width=Inches(5), height=Inches(3)
                    )
                    
                    overlay = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        self.prs.slide_width, self.prs.slide_height
                    )
                    overlay.fill.solid()
                    overlay.fill.fore_color.rgb = self.colors['dark']
                    overlay.fill.transparency = 0.85
                    overlay.line.fill.background()
                    
                    if picture:
                        overlay.zorder = 1
                        picture.zorder = 0
                
            elif slide_type == "section":
                slide = self.prs.slides.add_slide(self.section_header_layout)
                self.add_background_style(slide)
                
                title = slide.shapes.title
                title.text = slide_title
                self.apply_text_formatting(title, size=32, color=self.colors['primary'], bold=True, alignment=PP_ALIGN.CENTER)
                
                image_stream = self.get_image_from_unsplash(slide_title)
                if image_stream:
                    self.add_image_to_slide(
                        slide, image_stream,
                        left=Inches(3), top=Inches(2),
                        width=Inches(4), height=Inches(2.5)
                    )
                
                if slide_content:
                    content_box = slide.shapes.add_textbox(
                        Inches(2), Inches(4.5), Inches(6), Inches(0.75)
                    )
                    content_box.text_frame.text = "\n".join(slide_content[:1])
                    self.apply_text_formatting(content_box, size=16, color=self.colors['secondary'], alignment=PP_ALIGN.CENTER)
            
            else:
                slide = self.prs.slides.add_slide(self.title_content_layout)
                self.add_background_style(slide)
                self.create_slide_title(slide, slide_title, size=28)
                
                image_stream = None
                if i % 2 == 0:
                    image_stream = self.get_image_from_unsplash(slide_title)
                
                if slide_type == "comparison" and image_stream is None:
                    slide = self.prs.slides.add_slide(self.comparison_layout)
                    self.add_background_style(slide)
                    self.create_slide_title(slide, slide_title, size=28, alignment=PP_ALIGN.CENTER)
                    
                    left_content = []
                    right_content = []
                    column_mode = "left"
                    
                    for line in slide_content:
                        line = line.strip()
                        if not line:
                            continue
                            
                        if line.lower() in ["vs", "versus", "vs."]:
                            column_mode = "right"
                            continue
                        
                        if ":" in line and not left_content and not right_content:
                            parts = line.split(":")
                            left_content.append(parts[0])
                            if len(parts) > 1:
                                right_content.append(parts[1])
                            column_mode = "right"
                        else:
                            if column_mode == "left":
                                left_content.append(line)
                            else:
                                right_content.append(line)
                    
                    if not right_content:
                        midpoint = len(left_content) // 2
                        right_content = left_content[midpoint:]
                        left_content = left_content[:midpoint]
                    
                    left_box = slide.placeholders[1]
                    right_box = slide.placeholders[2]
                    
                    if left_content:
                        self.add_bullet_points(left_box, left_content)
                    
                    if right_content:
                        self.add_bullet_points(right_box, right_content)
                
                elif image_stream:
                    slide = self.prs.slides.add_slide(self.two_content_layout)
                    self.add_background_style(slide)
                    self.create_slide_title(slide, slide_title, size=28)
                    
                    content_shape = slide.placeholders[1]
                    self.add_bullet_points(content_shape, slide_content)
                    
                    self.add_image_to_slide(
                        slide, image_stream,
                        left=Inches(6), top=Inches(1.8),
                        width=Inches(3.5), height=Inches(2.5)
                    )
                else:
                    content_shape = slide.placeholders[1]
                    self.add_bullet_points(content_shape, slide_content)
            
            self.add_footer(slide, "Created with Marketing AI Assistant")
        
        if not any("conclusion" in slide["title"].lower() for slide in slides_data):
            closing_slide = self.prs.slides.add_slide(self.section_header_layout)
            self.add_background_style(closing_slide)
            
            title = closing_slide.shapes.title
            title.text = "Thank You"
            self.apply_text_formatting(title, size=36, color=self.colors['primary'], bold=True, alignment=PP_ALIGN.CENTER)
            
            contact_box = closing_slide.shapes.add_textbox(
                Inches(2), Inches(3), Inches(6), Inches(0.75)
            )
            contact_p = contact_box.text_frame.add_paragraph()
            contact_p.text = "Questions? Contact us today!"
            contact_p.alignment = PP_ALIGN.CENTER
            contact_p.font.size = Pt(20)
            contact_p.font.color.rgb = self.colors['secondary']
            
            logo_box = closing_slide.shapes.add_textbox(
                Inches(3), Inches(4), Inches(4), Inches(0.5)
            )
            logo_p = logo_box.text_frame.add_paragraph()
            logo_p.text = "Your Company"
            logo_p.alignment = PP_ALIGN.CENTER
            logo_p.font.size = Pt(16)
            logo_p.font.color.rgb = self.colors['accent1']
            logo_p.font.bold = True
        
        output_path = 'static/generated_ppt.pptx'
        self.prs.save(output_path)
        return output_path

# Add login route
@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')
        
        if username == 'sanath' and password == 'Sanath@123':
            session['logged_in'] = True
            return redirect(get_redirect_url('index'))
        else:
            flash('Invalid credentials. Please try again.', 'error')
    
    # If user is already logged in, redirect to index
    if session.get('logged_in'):
        return redirect(get_redirect_url('index'))
            
    return render_template('login.html')

# Add logout route
@app.route('/logout')
def logout():
    session.pop('logged_in', None)
    return redirect(url_for('login'))

# Modify index route to check for login
@app.route('/')
def index():
    if not session.get('logged_in'):
        return redirect(url_for('login'))
    return render_template('index.html')

# Add an explicit route for /marketingai/ to handle redirections after login
@app.route('/marketingai/')
def marketingai_index():
    if not session.get('logged_in'):
        return redirect(url_for('login'))
    return render_template('index.html')

# Make sure all API routes are protected
def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if not session.get('logged_in'):
            return jsonify({'error': 'Not authorized'}), 401
        return f(*args, **kwargs)
    return decorated_function

@app.route('/upload', methods=['POST'])
@login_required
def upload_file():
    try:
        print("Upload endpoint called")
        if 'file' not in request.files:
            print("No file provided in request")
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            print("No file selected")
            return jsonify({'error': 'No file selected'}), 400
        
        print(f"Processing file: {file.filename}")
        
        if file and file.filename.lower().endswith('.pdf'):
            if rag_engine.is_file_processed(file.filename):
                print(f"File '{file.filename}' already processed")
                return jsonify({
                    'message': 'File already processed',
                    'loaded_files': rag_engine.loaded_files  # Changed from get_loaded_files() to loaded_files
                })
            
            filepath = os.path.join(UPLOAD_FOLDER, file.filename)
            print(f"Saving file to {filepath}")
            os.makedirs(os.path.dirname(filepath), exist_ok=True)
            file.save(filepath)
            
            if not os.path.exists(filepath):
                print(f"File failed to save at {filepath}")
                return jsonify({'error': 'Failed to save uploaded file'}), 500
                
            print(f"File saved, size: {os.path.getsize(filepath)} bytes")
            
            print("Starting PDF processing")
            texts, documents = pdf_processor.process_pdf(filepath)
            
            print(f"PDF processing completed: {len(texts)} texts, {len(documents)} documents")
            
            if not texts or len(texts) == 0:
                print("No text could be extracted from the PDF")
                return jsonify({'error': 'No text could be extracted from the PDF. The file might be scanned or contain only images.'}), 400
            
            print(f"Adding to index: {file.filename}")
            success = rag_engine.add_to_index(texts, documents, file.filename)
            
            if not success:
                print("Failed to add to vector index")
                return jsonify({'error': 'Failed to add document to vector index'}), 500
                
            print(f"Successfully processed file {file.filename}")
            
            return jsonify({
                'message': 'PDF processed successfully',
                'chunks': len(texts),
                'loaded_files': rag_engine.loaded_files  # Changed from get_loaded_files() to loaded_files
            })
        
        print(f"Invalid file type: {file.filename}")
        return jsonify({'error': 'Invalid file type. Please upload PDF files only.'}), 400
        
    except Exception as e:
        print(f"Error processing upload: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': f'Processing error: {str(e)}'}), 500

@app.route('/files', methods=['GET'])
@login_required
def list_files():
    # Change this line to use the loaded_files attribute directly
    return jsonify({'files': rag_engine.loaded_files})

@app.route('/chat', methods=['POST'])
@login_required
def chat():
    try:
        data = request.json
        query = data.get('query')
        
        if not query:
            return jsonify({'error': 'No query provided'}), 400
        
        content_params = data.get('content_params', {})
        
        # Get or create session ID for memory management
        session_id = session.get('chat_session_id')
        if not session_id:
            session_id = f"session_{int(time.time())}_{os.urandom(4).hex()}"
            session['chat_session_id'] = session_id
        
        # Get response with session-based memory
        response, _ = rag_engine.get_response(query, content_params, session_id)
        
        return jsonify({
            'response': response,
            'session_id': session_id
        }), 200
        
    except Exception as e:
        print(f"Error in chat endpoint: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({
            'error': 'An error occurred processing your request',
            'details': str(e)
        }), 500

@app.route('/new-chat', methods=['POST'])
@login_required
def new_chat():
    """Start a new chat session by clearing memory"""
    try:
        # Get current session ID
        current_session_id = session.get('chat_session_id')
        
        # Clear memory for current session if it exists
        if current_session_id:
            rag_engine.clear_session_memory(current_session_id)
        
        # Create new session ID
        new_session_id = f"session_{int(time.time())}_{os.urandom(4).hex()}"
        session['chat_session_id'] = new_session_id
        
        return jsonify({
            'message': 'New chat session started',
            'session_id': new_session_id
        }), 200
        
    except Exception as e:
        print(f"Error starting new chat: {str(e)}")
        return jsonify({'error': f'Error starting new chat: {str(e)}'}), 500

@app.route('/chat-history', methods=['GET'])
@login_required
def get_chat_history():
    """Get conversation history for current session"""
    try:
        session_id = session.get('chat_session_id', 'default')
        history = rag_engine.get_session_history(session_id)
        
        # Convert messages to a serializable format
        formatted_history = []
        for message in history:
            if hasattr(message, 'content'):
                role = 'user' if message.__class__.__name__ == 'HumanMessage' else 'assistant'
                formatted_history.append({
                    'role': role,
                    'content': message.content,
                    'timestamp': getattr(message, 'timestamp', None)
                })
        
        return jsonify({
            'history': formatted_history,
            'session_id': session_id
        }), 200
        
    except Exception as e:
        print(f"Error getting chat history: {str(e)}")
        return jsonify({'error': f'Error getting chat history: {str(e)}'}), 500

@app.route('/clear-memory', methods=['POST'])
@login_required
def clear_memory():
    """Clear conversation memory for current session"""
    try:
        session_id = session.get('chat_session_id', 'default')
        rag_engine.clear_session_memory(session_id)
        
        return jsonify({
            'message': 'Conversation memory cleared',
            'session_id': session_id
        }), 200
        
    except Exception as e:
        print(f"Error clearing memory: {str(e)}")
        return jsonify({'error': f'Error clearing memory: {str(e)}'}), 500

@app.route('/generate-ppt', methods=['POST'])
@login_required
def generate_ppt():
    try:
        data = request.json
        content = data.get('content')
        
        if not content:
            return jsonify({'error': 'No content provided'}), 400
            
        content = content.replace('[PPT_CONTENT_START]', '').replace('[PPT_CONTENT_END]', '')
        
        # Generate a unique filename with timestamp to prevent caching issues
        timestamp = int(time.time())
        filename = f'generated_ppt_{timestamp}.pptx'
        output_path = os.path.join('static', filename)
        
        generator = PPTGenerator()
        generator.create_presentation(content)
        generator.prs.save(output_path)
        
        # Log the file creation
        print(f"PPT file created: {output_path}")
        
        # Create a direct download endpoint URL
        download_url = url_for('download_ppt', filename=filename)
        
        # Make sure the URL includes APPLICATION_ROOT if defined
        if APPLICATION_ROOT and not download_url.startswith(APPLICATION_ROOT):
            download_url = APPLICATION_ROOT + download_url
        
        # Log the download URL for debugging
        print(f"Download URL: {download_url}")
            
        return jsonify({
            'message': 'Presentation generated successfully',
            'download_url': download_url
        })
        
    except Exception as e:
        import traceback
        print(f"Error generating PPT: {str(e)}")
        print(traceback.format_exc())
        return jsonify({'error': f'Error generating PPT: {str(e)}'}), 500

@app.route('/download-ppt/<filename>')
@login_required
def download_ppt(filename):
    try:
        file_path = os.path.join(app.root_path, 'static', filename)
        
        if not os.path.exists(file_path):
            return jsonify({'error': 'File not found'}), 404
            
        # Use Flask's send_file with the correct MIME type and attachment filename
        return send_file(
            file_path,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name='presentation.pptx'
        )
    except Exception as e:
        print(f"Error downloading PPT: {str(e)}")
        return jsonify({'error': f'Error downloading file: {str(e)}'}), 500

@app.route('/delete-file', methods=['POST'])
@login_required
def delete_file():
    try:
        data = request.json
        filename = data.get('filename')
        
        if not filename:
            return jsonify({'error': 'No filename provided'}), 400
        
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        
        try:
            if os.path.exists(filepath):
                os.remove(filepath)
                print(f"Physical file deleted: {filepath}")
        except Exception as e:
            print(f"Error deleting physical file: {e}")
        
        if rag_engine.remove_file(filename):
            print(f"File removed from index: {filename}")
            return jsonify({'message': 'File deleted successfully'})
        else:
            return jsonify({'error': 'File not found in index'}), 404
            
    except Exception as e:
        print(f"Error in delete_file: {str(e)}")
        return jsonify({'error': f'Error deleting file: {str(e)}'}), 500

@app.route('/suggest-content', methods=['POST'])
@login_required
def suggest_content():
    try:
        data = request.json
        topic = data.get('topic')
        content_type = data.get('content_type')
        
        if not topic:
            return jsonify({'error': 'No topic provided'}), 400
            
        # Adapt this to generate academic content suggestions instead of marketing
        prompt = f"Suggest 5 study resources or academic content ideas about {topic}"
        if content_type:
            prompt += f" for {content_type} study material"
            
        suggestions = rag_engine.get_response(prompt, None)[0]
        
        return jsonify({
            'suggestions': suggestions
        })
        
    except Exception as e:
        return jsonify({'error': f'Error generating suggestions: {str(e)}'}), 500

@app.route('/generate-title', methods=['POST'])
@login_required
def generate_title():
    try:
        data = request.json
        message = data.get('message', '')
        
        if not message:
            return jsonify({'title': 'New Chat'}), 200
        
        # Extract first 5-7 words for a concise title
        words = message.split()[:6]
        title = ' '.join(words)
        if len(words) < len(message.split()):
            title += '...'
            
        return jsonify({'title': title})
        
    except Exception as e:
        print(f"Error generating title: {str(e)}")
        return jsonify({'error': str(e), 'title': 'New Chat'}), 500

@app.route('/generate-mcq', methods=['POST'])
@login_required
def generate_mcq():
    try:
        data = request.json
        subject = data.get('subject', '')
        num_questions = int(data.get('num_questions', 5))
        
        if not subject:
            return jsonify({'error': 'No subject provided'}), 400
            
        # Use the RAG engine to generate MCQ test
        mcq_test = rag_engine.generate_mcq_test(subject, num_questions)
        
        if not mcq_test:
            return jsonify({'error': 'Failed to generate MCQ test'}), 500
            
        return jsonify({
            'success': True,
            'mcq_test': mcq_test
        })
        
    except Exception as e:
        print(f"Error generating MCQ: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': f'Error generating MCQ: {str(e)}'}), 500

@app.route('/evaluate-mcq', methods=['POST'])
@login_required
def evaluate_mcq():
    try:
        data = request.json
        mcq_test = data.get('mcq_test', [])
        user_answers = data.get('user_answers', {})
        
        if not mcq_test or not user_answers:
            return jsonify({'error': 'Missing test or answers'}), 400
            
        # Use the RAG engine to evaluate the answers
        results = rag_engine.evaluate_mcq_answers(mcq_test, user_answers)
        
        return jsonify({
            'success': True,
            'results': results
        })
        
    except Exception as e:
        print(f"Error evaluating MCQ: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': f'Error evaluating MCQ: {str(e)}'}), 500

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=8502)